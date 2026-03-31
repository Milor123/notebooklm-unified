#!/usr/bin/env python3
"""
NotebookLM Workflow - Unificado
Une PPTXs -> PDF sin pérdida -> Elimina watermark -> PPTX limpio

Uso: python workflow.py [opciones]

Ejemplo: python workflow.py --input "carpeta_con_pptx" --output "resultado.pptx"
"""

import argparse
import os
import sys
import subprocess
import tempfile
import shutil
from pathlib import Path
from io import BytesIO
import re
import xml.etree.ElementTree as ET
import zipfile


# ============================================================
# AUTO-INSTALACIÓN DE DEPENDENCIAS
# ============================================================


def instalar_dependencias():
    """Instala las dependencias necesarias si no están disponibles."""
    paquetes = ["python-pptx", "img2pdf", "PyMuPDF", "Pillow"]

    print("🔍 Verificando dependencias...")
    faltantes = []

    for pkg in paquetes:
        # Usar pip show para verificar de forma más precisa
        try:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "show", pkg],
                capture_output=True,
                text=True,
                timeout=15
            )
            if result.returncode == 0:
                print(f"  ✓ {pkg} ya instalado")
            else:
                print(f"  ✗ {pkg} no encontrado, se instalará...")
                faltantes.append(pkg)
        except Exception as e:
            print(f"  ✗ {pkg}: error al verificar, se instalará...")
            faltantes.append(pkg)

    if faltantes:
        print(f"\n📦 Instalando dependencias: {', '.join(faltantes)}")
        try:
            subprocess.check_call(
                [sys.executable, "-m", "pip", "install"] + faltantes + ["--quiet"]
            )
            print("✅ Dependencias instaladas correctamente")

            # Verificar que realmente se instalaron
            print("🔍 Verificando instalación...")
            for pkg in faltantes:
                result = subprocess.run(
                    [sys.executable, "-m", "pip", "show", pkg],
                    capture_output=True,
                    text=True,
                    timeout=15
                )
                if result.returncode == 0:
                    print(f"  ✓ {pkg} instalado correctamente")
                else:
                    print(f"  ⚠ {pkg} puede no haber instalado correctamente")

        except subprocess.CalledProcessError as e:
            print(f"❌ Error al instalar dependencias: {e}")
            sys.exit(1)


# ============================================================
# CLASE PRINCIPAL DEL WORKFLOW
# ============================================================


class NotebookLMWorkflow:
    """Workflow unificado para procesar presentaciones de NotebookLM."""

    def __init__(self, verbose=False, debug=False):
        self.verbose = verbose
        self.debug = debug
        self.temp_dir = None
        self.directorio_trabajo = None

    def log(self, mensaje):
        if self.verbose:
            print(f"  → {mensaje}")

    def iniciar(self):
        """Inicia el directorio temporal de trabajo."""
        self.temp_dir = tempfile.mkdtemp(prefix="notebooklm_workflow_")
        self.directorio_trabajo = Path(self.temp_dir)
        print(f"📁 Directorio temporal: {self.directorio_trabajo}")

    def limpiar(self):
        """Limpia los archivos temporales."""
        if self.temp_dir and os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
            print("🧹 Directorio temporal limpiado")

    # ----------------------------------------------------------
    # PASO 1: UNIR PPTXs
    # ----------------------------------------------------------

    def unir_ppts(self, archivos_pptx):
        """Une múltiples archivos PPTX en uno solo."""
        from pptx import Presentation

        if not archivos_pptx:
            raise ValueError("No se proporcionaron archivos PPTX")

        print(f"\n📑 UNIENDO {len(archivos_pptx)} archivos PPTX...")

        # Ordenar por nombre
        archivos_ordenados = sorted(archivos_pptx, key=lambda p: p.name.lower())

        for f in archivos_ordenados:
            print(f"  - {f.name}")

        # Cargar primer archivo para obtener dimensiones
        prs_base = Presentation(str(archivos_ordenados[0]))
        prs_final = Presentation()
        prs_final.slide_width = prs_base.slide_width
        prs_final.slide_height = prs_base.slide_height

        # Eliminar slide inicial vacío
        self._eliminar_slide_inicial(prs_final)

        total_slides = 0

        for archivo in archivos_ordenados:
            self.log(f"Procesando: {archivo.name}")
            prs = Presentation(str(archivo))

            for slide in prs.slides:
                self._copiar_slide_solo_imagenes(prs_final, slide)
                total_slides += 1

        # Guardar PPTX unido
        pptx_unido = self.directorio_trabajo / "unido.pptx"
        prs_final.save(str(pptx_unido))

        print(f"  ✓ PPTX unite con {total_slides} diapositivas")
        return pptx_unido

    def _eliminar_slide_inicial(self, prs):
        """Elimina el slide inicial vacío del PPTX."""
        if len(prs.slides) == 0:
            return
        slide_id = prs.slides._sldIdLst[0]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]

    def _copiar_slide_solo_imagenes(self, prs_destino, slide_origen):
        """Copia solo las imágenes de un slide al destino."""
        nueva = prs_destino.slides.add_slide(prs_destino.slide_layouts[6])

        for shape in slide_origen.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                img_bytes = shape.image.blob
                nueva.shapes.add_picture(
                    BytesIO(img_bytes),
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )

    # ----------------------------------------------------------
    # PASO 2: CONVERTIR PPTX A PDF (LOSSLESS)
    # ----------------------------------------------------------

    def pptx_a_pdf(self, pptx_path):
        """Convierte PPTX a PDF sin pérdida de calidad."""
        import img2pdf

        print("\n🔄 CONVIERTIENDO A PDF (sin pérdida)...")

        # Extraer imágenes del PPTX
        imagenes = self._extraer_imagenes_de_pptx(pptx_path)

        if not imagenes:
            raise ValueError("No se pudieron extraer imágenes del PPTX")

        print(f"  ✓ Extraídas {len(imagenes)} imágenes")

        # Convertir a PDF usando img2pdf
        pdf_path = pptx_path.with_suffix(".pdf")

        with open(pdf_path, "wb") as f:
            f.write(img2pdf.convert([str(p) for p in imagenes]))

        print(f"  ✓ PDF creado: {pdf_path.name}")
        return pdf_path

    def _extraer_imagenes_de_pptx(self, pptx_path):
        """Extrae imágenes de un PPTX en orden."""
        NS = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
        }

        IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".gif", ".webp"}
        imagenes = []

        with zipfile.ZipFile(pptx_path, "r") as zf:
            # Obtener slides en orden
            slides = []
            for name in zf.namelist():
                m = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", name)
                if m:
                    slides.append((int(m.group(1)), name))
            slides.sort()
            slide_files = [name for _, name in slides]

            for idx, slide_xml in enumerate(slide_files, start=1):
                images = self._get_slide_images_with_size(zf, slide_xml, NS, IMAGE_EXTS)

                if not images:
                    print(f"  [WARN] Slide {idx}: no se encontraron imágenes")
                    continue

                # Elegir la imagen más grande
                img_path_in_zip, area = max(images, key=lambda x: x[1])

                if img_path_in_zip not in zf.namelist():
                    raise FileNotFoundError(f"No existe en PPTX: {img_path_in_zip}")

                ext = Path(img_path_in_zip).suffix.lower()
                out_path = self.directorio_trabajo / f"slide_{idx:04d}{ext}"

                with zf.open(img_path_in_zip) as src, open(out_path, "wb") as dst:
                    shutil.copyfileobj(src, dst)

                imagenes.append(out_path)

        return imagenes

    def _get_slide_images_with_size(self, zf, slide_xml_path, NS, IMAGE_EXTS):
        """Obtiene las imágenes de un slide con su tamaño."""
        # Obtener mapa de relaciones
        relmap = self._get_relationship_map(zf, slide_xml_path, NS)

        # Parsear XML del slide
        slide_root = ET.fromstring(zf.read(slide_xml_path))

        out = []

        for pic in slide_root.findall(".//p:pic", NS):
            blip = pic.find(".//a:blip", NS)
            if blip is None:
                continue

            rid = blip.attrib.get(f"{{{NS['r']}}}embed")
            if not rid:
                continue

            target = relmap.get(rid)
            if not target:
                continue

            if Path(target).suffix.lower() not in IMAGE_EXTS:
                continue

            # Obtener tamaño
            xfrm = pic.find(".//a:xfrm", NS)
            cx = cy = 0
            if xfrm is not None:
                ext = xfrm.find("a:ext", NS)
                if ext is not None:
                    cx = int(ext.attrib.get("cx", "0"))
                    cy = int(ext.attrib.get("cy", "0"))

            area = cx * cy
            out.append((target, area))

        return out

    def _get_relationship_map(self, zf, slide_xml_path, NS):
        """Obtiene el mapa de relaciones de un slide."""
        from pathlib import PurePosixPath

        rels_path = slide_xml_path.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
        relmap = {}

        if rels_path not in zf.namelist():
            return relmap

        rels_root = ET.fromstring(zf.read(rels_path))

        def normalize_zip_path(path_str):
            parts = []
            for part in PurePosixPath(path_str).parts:
                if part in ("", "."):
                    continue
                if part == "..":
                    if parts:
                        parts.pop()
                else:
                    parts.append(part)
            return "/".join(parts)

        def resolve_target(target):
            base_dir = PurePosixPath(slide_xml_path).parent
            combined = str(base_dir / target)
            return normalize_zip_path(combined)

        for rel in rels_root.findall("rel:Relationship", NS):
            rid = rel.attrib.get("Id")
            target = rel.attrib.get("Target")
            if rid and target:
                relmap[rid] = resolve_target(target)

        return relmap

    # ----------------------------------------------------------
    # PASO 3: ELIMINAR WATERMARK (del app.py original)
    # ----------------------------------------------------------

    def eliminar_watermark(self, pdf_path, debug=False):
        """Elimina el watermark de NotebookLM del PDF."""
        import fitz
        from PIL import Image
        import io

        print("\n🧹 ELIMINANDO WATERMARK...")

        doc = fitz.open(str(pdf_path))
        pages_processed = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            rect = page.rect

            # NotebookLM watermark posición -角落右下
            # Ajustable según necesidad
            wm_x1 = rect.width - 115
            wm_y1 = rect.height - 30
            wm_x2 = rect.width - 5
            wm_y2 = rect.height - 5

            if debug:
                print(f"  Página {page_num + 1}: size={rect.width}x{rect.height}, watermark zone=({wm_x1},{wm_y1})-({wm_x2},{wm_y2})")

            # Muestrear color de fondo acima del watermark
            # Verificar que el rectángulo de muestreo sea válido
            if wm_y1 - 10 > 0:
                sample_rect = fitz.Rect(wm_x1, wm_y1 - 10, wm_x2, wm_y1 - 2)
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat, clip=sample_rect)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                pixels = img.load()

                # Dibujar columna por columna usando el color de fondo
                col_width = (wm_x2 - wm_x1) / img.width
                for x in range(img.width):
                    color = pixels[x, img.height // 2]
                    r = color[0] / 255
                    g = color[1] / 255
                    b = color[2] / 255

                    col_rect = fitz.Rect(
                        wm_x1 + x * col_width, wm_y1, wm_x1 + (x + 1) * col_width, wm_y2
                    )
                    page.draw_rect(col_rect, color=(r, g, b), fill=(r, g, b))
            else:
                print(f"  [WARN] Página {page_num + 1}: no se puede muestrear (altura muy pequeña)")

            pages_processed += 1

        # Guardar PDF limpio
        pdf_limpio = pdf_path.parent / f"cleaned_{pdf_path.name}"
        doc.save(pdf_limpio, garbage=4, deflate=True, clean=True)
        doc.close()

        print(f"  ✓ {pages_processed} páginas procesadas")
        return pdf_limpio

    # ----------------------------------------------------------
    # PASO 4: PDF A PPTX
    # ----------------------------------------------------------

    def pdf_a_pptx(self, pdf_path, nombre_salida="resultado.pptx"):
        """Convierte PDF limpio de vuelta a PPTX."""
        import fitz

        print("\n📊 CONVIERTIENDO A PPTX...")

        doc = fitz.open(str(pdf_path))

        if len(doc) == 0:
            raise ValueError("PDF vacío")

        from pptx import Presentation

        # Usar tamaño de primera página
        page0 = doc[0]
        rect = page0.rect
        base_width = rect.width
        base_height = rect.height

        prs = Presentation()

        # Mantener proporción
        ancho_pulgadas = 10
        alto_pulgadas = ancho_pulgadas * (base_height / base_width)

        prs.slide_width = int(ancho_pulgadas * 914400)
        prs.slide_height = int(alto_pulgadas * 914400)

        # Eliminar slide inicial vacío
        self._eliminar_slide_inicial(prs)

        total = len(doc)
        print(f"  → {total} páginas a procesar")

        for i, page in enumerate(doc, start=1):
            # Alta calidad sin compresión
            pix = page.get_pixmap(dpi=200, alpha=False)
            img_bytes = pix.tobytes("png")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                BytesIO(img_bytes),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

            print(f"  Página {i}/{total}")

        # Guardar
        pptx_salida = pdf_path.parent / nombre_salida
        prs.save(str(pptx_salida))

        print(f"  ✓ PPTX creado: {pptx_salida.name}")
        return pptx_salida

    # ----------------------------------------------------------
    # EJECUTAR WORKFLOW COMPLETO
    # ----------------------------------------------------------

    def ejecutar(self, archivos_pptx, archivo_salida=None):
        """Ejecuta el workflow completo."""
        try:
            # Iniciar
            self.iniciar()

            # Paso 1: Unir PPTXs
            pptx_unido = self.unir_ppts(archivos_pptx)

            # Paso 2: Convertir a PDF
            pdf = self.pptx_a_pdf(pptx_unido)

            # Paso 3: Eliminar watermark
            pdf_limpio = self.eliminar_watermark(pdf, debug=self.debug)

            # Paso 4: Convertir a PPTX
            nombre_salida = archivo_salida or "presentacion_limpia.pptx"
            pptx_final = self.pdf_a_pptx(pdf_limpio, nombre_salida)

            print("\n" + "="*50)
            print("✅ WORKFLOW COMPLETADO")
            print("="*50)

            print(f"\n📦 Archivo final: {pptx_final.absolute()}")

            # Copiar a directorio actual
            destino_final = Path.cwd() / pptx_final.name
            shutil.copy2(pptx_final, destino_final)
            print(f"📦 Copiado a: {destino_final}")

            return destino_final

        finally:
            # Pequeña espera para liberar archivos en Windows
            import time
            time.sleep(0.5)
            self.limpiar()


# ============================================================
# INTERFAZ CLI
# ============================================================


def main():
    parser = argparse.ArgumentParser(
        description="Workflow unificado para limpiar presentaciones de NotebookLM",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos de uso:
  python workflow.py                                    # Une todos los .pptx del directorio actual
  python workflow.py --input carpeta                    # Desde carpeta específica
  python workflow.py --output mi_presentacion.pptx     # Nombre de salida personalizado
  python workflow.py -v                                # Modo verboso
  python workflow.py -d                                # Debug: muestra coordenadas del watermark
        """,
    )

    parser.add_argument(
        "--input",
        "-i",
        help="Carpeta o archivo(s) PPTX de entrada. Si es carpeta, usa todos los .pptx",
        default=".",
    )

    parser.add_argument(
        "--output",
        "-o",
        help="Nombre del archivo PPTX de salida",
        default="presentacion_limpia.pptx",
    )

    parser.add_argument("--verbose", "-v", action="store_true", help="Modo verboso")
    parser.add_argument(
        "--debug",
        "-d",
        action="store_true",
        help="Debug: muestra coordenadas del watermark",
    )

    args = parser.parse_args()

    # Instalar dependencias primero
    instalar_dependencias()

    # Determinar archivos de entrada
    input_path = Path(args.input)

    if input_path.is_file():
        archivos = [input_path]
    elif input_path.is_dir():
        archivos = sorted(input_path.glob("*.pptx"), key=lambda p: p.name.lower())
        # Filtrar archivos unidos anteriores y limpiados
        archivos = [
            a
            for a in archivos
            if a.name.lower() not in {"unido.pptx", "unido_bien.pptx", "unido.pptx"}
            and not a.name.lower().startswith("cleaned_")
        ]
    else:
        # Buscar en directorio actual
        archivos = sorted(Path(".").glob("*.pptx"), key=lambda p: p.name.lower())
        archivos = [
            a
            for a in archivos
            if a.name.lower() not in {"unido.pptx", "unido_bien.pptx", "unido.pptx"}
            and not a.name.lower().startswith("cleaned_")
        ]

    if not archivos:
        print("❌ No se encontraron archivos PPTX")
        print("Usage: python workflow.py --input <carpeta>")
        sys.exit(1)

    print("=" * 50)
    print("🚀 NotebookLM Workflow - Inicio")
    print("=" * 50)

    # Ejecutar workflow
    workflow = NotebookLMWorkflow(verbose=args.verbose, debug=args.debug)
    resultado = workflow.ejecutar(archivos, args.output)

    print(f"\n🎉 Listo! Archivo disponible: {resultado}")


if __name__ == "__main__":
    main()
